"""
Gerador Programático de Apresentações PPTX (python-pptx).

Cria automaticamente os slides de estudo para reuniões de pequeno grupo / células,
escola bíblica e cultos de ensino.
"""

import os
import logging
from typing import Dict, Any, List, Optional
from pathlib import Path

import sys
sys.path.append(str(Path(__file__).resolve().parent.parent.parent))
from config.settings import get_folder_path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s")
logger = logging.getLogger(__name__)


class CellSlidePPTXGenerator:
    """
    Gerador de apresentações PowerPoint para células e estudos bíblicos.
    """

    def __init__(self):
        """
        Inicializa a pasta de apresentações.
        """
        self.output_dir = get_folder_path("SLIDES_ESTUDO")
        os.makedirs(self.output_dir, exist_ok=True)

    def generate_cell_slides(
        self,
        study_title: str,
        biblical_passage: str,
        discussion_points: List[str],
        output_filename: str = "estudo_celula_semanal.pptx"
    ) -> str:
        """
        Gera o arquivo .pptx para o estudo de célula da semana.

        :param study_title: Título do Estudo.
        :param biblical_passage: Texto Bíblico Base.
        :param discussion_points: Perguntas / Pontos de Aplicação Prática.
        :param output_filename: Nome do arquivo PPTX.
        :return: Caminho do arquivo gerado.
        """
        output_path = os.path.join(self.output_dir, output_filename)
        logger.info(f"📊 Gerando slides PPTX para o Estudo de Célula: '{study_title}'...")

        if not HAS_PPTX:
            return self._mock_pptx_file(output_path)

        try:
            prs = Presentation()

            # Slide 1: Capa
            slide_layout = prs.slide_layouts[0]
            slide1 = prs.slides.add_slide(slide_layout)
            title1 = slide1.shapes.title
            subtitle1 = slide1.placeholders[1]

            title1.text = study_title
            subtitle1.text = f"Estudo de Célula IBPM CR\nTexto Base: {biblical_passage}"

            # Slide 2: Quebra-gelo e Leitura
            slide2 = prs.slides.add_slide(prs.slide_layouts[1])
            slide2.shapes.title.text = "1. Leitura e Quebra-Gelo"
            tf2 = slide2.placeholders[1].text_frame
            tf2.text = f"Vamos ler juntos a passagem de {biblical_passage}."
            p = tf2.add_paragraph()
            p.text = "Pergunta inicial: Como este ensino se aplica à nossa rotina durante esta semana?"

            # Slide 3: Perguntas de Aplicação Prática
            slide3 = prs.slides.add_slide(prs.slide_layouts[1])
            slide3.shapes.title.text = "2. Pontos para Discussão em Célula"
            tf3 = slide3.placeholders[1].text_frame
            for pt in discussion_points:
                p = tf3.add_paragraph()
                p.text = f"• {pt}"

            prs.save(output_path)
            logger.info(f"✅ Apresentação PPTX salva em: {output_path}")
            return output_path

        except Exception as e:
            logger.error(f"❌ Erro ao gerar apresentação PPTX: {e}")
            return self._mock_pptx_file(output_path)

    def _mock_pptx_file(self, output_path: str) -> str:
        """Gera um arquivo PPTX placeholder."""
        with open(output_path, "wb") as f:
            f.write(b"MOCK_PPTX_PRESENTATION_DATA")
        logger.info(f"📁 Slide PPTX gerado (placeholder): {output_path}")
        return output_path


if __name__ == "__main__":
    pptx_gen = CellSlidePPTXGenerator()
    points = [
        "O que significa apresentar o corpo como sacrifício vivo?",
        "De que maneira o mundo tenta nos conformar com suas práticas?",
        "Qual foi a experiência de transformação de mente que você já viveu?"
    ]
    out = pptx_gen.generate_cell_slides("Vencendo pela Palavra", "Romanos 12:1-2", points)
    print(f"Apresentação gerada em: {out}")
